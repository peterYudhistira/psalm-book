from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from pptx.slide import Slide
import copy
import os

class LyricParams:
    def __init__(self, lyricLine1, lyricLine2, translationLine1, translationLine2, fontName, fontSize):
        self.lyricLine1 = lyricLine1
        self.lyricLine2= lyricLine2
        self.translationLine1 = translationLine1
        self.translationLine2 = translationLine2
        self.fontName = fontName
        self.fontSize = fontSize

class PPTHandler:
    def __init__(self, pptFilePath:str, lyricParams:LyricParams):
        self.path = pptFilePath
        self.presentation = Presentation(pptFilePath)
        self.width = self.presentation.slide_width
        self.height = self.presentation.slide_height
        self.lyricParams = lyricParams

        # open the ppt file


    # copies from example.pptx (is a fixed file)
    def CopySlide(self, slideNumber:int):

        # source is example.pptx
        src = Presentation("example.pptx")
        
        # specify the slide you want to copy the contents from
        src_slide = src.slides[slideNumber - 1]

        # Define the layout you want to use from your generated pptx
        slide_layout = None
        for layout in self.presentation.slide_layouts:
            if "blank" in layout.name.lower():
                slide_layout = layout
                break
            
        # no validation, we die like men
        # create now slide, to copy contents to
        curr_slide = self.presentation.slides.add_slide(slide_layout)

        # create images dict
        imgDict = {}

        # create textbox dict
        textBoxList = [] 

        # now copy contents from external slide, but do not copy slide properties
        # e.g. slide layouts, etc., because these would produce errors, as diplicate
        # entries might be generated
        for shp in src_slide.shapes:
            if 'Picture' in shp.name:
                # save image
                with open(shp.name+'.jpg', 'wb') as f:
                    f.write(shp.image.blob)

                # add image to dict
                imgDict[shp.name+'.jpg'] = [shp.left, shp.top, shp.width, shp.height]

            elif 'TextBox' in shp.name:
                textBoxDict = {}
                textBoxDict['left'] = shp.left
                textBoxDict['top'] = shp.top
                textBoxDict['width'] = shp.width
                textBoxDict['height'] = shp.height
                textBoxDict['text'] = shp.text

                textBoxList.append(textBoxDict)
                
        


        # add pictures
        for k, v in imgDict.items():
            curr_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
            os.remove(k)
        
        # add text
        for textBoxDict in textBoxList:
            new_shape = curr_slide.shapes.add_textbox(
                    left=textBoxDict['left'],
                    top=textBoxDict['top'], 
                    width=textBoxDict['width'],
                    height=textBoxDict['height']
                )
            new_shape.text = textBoxDict['text']

        # save the thing
        self.presentation.save(self.path)

    # saves into a directory
    def SaveToPath(self, savePath):
        pass

    # edits a slide based on the lyrics.
    def InputLyrics(self, slideNumber:int):
        # get current slide
        currentSlideNumber = slideNumber
        slide = self.presentation.slides[currentSlideNumber - 1]  # Index is 0-based, so subtract 1
        
        # center it
        lyricFrame = slide.shapes[1]
        lyricFrame.left = int((self.width - lyricFrame.width) / 2)
        lyricFrame.top = int((self.height - lyricFrame.height) / 2)

        # clear
        lyricTextBox = slide.shapes[1].text_frame
        lyricTextBox.clear()
        paragraph = lyricTextBox.paragraphs[0]

        # lyric 1
        lyricLine1 = paragraph.add_run()
        self.StyleLyrics(lyricLine1, self.lyricParams.lyricLine1)

        # tl 1
        translationLine1 = paragraph.add_run()
        self.StyleTranslation(translationLine1, self.lyricParams.translationLine1)
        
        # lyric 2
        lyricLine2 = paragraph.add_run()
        self.StyleLyrics(lyricLine2, self.lyricParams.lyricLine2)

        # tl 2
        translationLine2 = paragraph.add_run()
        self.StyleTranslation(translationLine2, self.lyricParams.translationLine2)
        

    def StyleLyrics(self, lyricLine, text):
        lyricLine.text = text
        lyricLine.font.bold = True
        lyricLine.font.name = self.lyricParams.fontName
        lyricLine.font.size = Pt(self.lyricParams.fontSize)

    def StyleTranslation(self, translationLine, text):
        translationLine.text = text
        translationLine.font.italic = True
        translationLine.font.name = self.lyricParams.fontName
        translationLine.font.size = Pt(self.lyricParams.fontSize)


lp = LyricParams(
    "haha",
    "hehe",
    "haha",
    "hehe",
    "Arial",
    32
)

ph = PPTHandler(
    "example2.pptx",
    lp
)

ph.CopySlide(2)
ph.InputLyrics(4)